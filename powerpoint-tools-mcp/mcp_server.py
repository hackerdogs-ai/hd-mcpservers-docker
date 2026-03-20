#!/usr/bin/env python3
"""PowerPoint Tools MCP Server — create and read PPTX presentations."""
import json
import logging
import os
import sys
from typing import Optional

from fastmcp import FastMCP

logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(name)s] %(levelname)s: %(message)s", stream=sys.stderr)
logger = logging.getLogger("powerpoint-tools-mcp")
MCP_TRANSPORT = os.environ.get("MCP_TRANSPORT", "stdio")
MCP_PORT = int(os.environ.get("MCP_PORT", "8507"))
mcp = FastMCP("PowerPoint Tools MCP Server", instructions="Create and read PPTX presentations using python-pptx.")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


@mcp.tool()
def create_presentation(slides_json: str, output_path: str, title: Optional[str] = None) -> str:
    """Create a PowerPoint presentation from JSON slide definitions.

    Args:
        slides_json: JSON array of slide objects. Each: {"title": "...", "content": "...", "layout": "title_and_content"}.
                     layout can be: title, title_and_content, section_header, blank.
        output_path: Output .pptx file path.
        title: Optional presentation title for the first slide.
    """
    if not PPTX_AVAILABLE:
        return json.dumps({"error": "python-pptx not installed"})
    try:
        prs = Presentation()
        slides = json.loads(slides_json)

        if title:
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title

        for s in slides:
            layout_name = s.get("layout", "title_and_content")
            layout_idx = {"title": 0, "title_and_content": 1, "section_header": 2, "blank": 6}.get(layout_name, 1)
            slide_layout = prs.slide_layouts[min(layout_idx, len(prs.slide_layouts) - 1)]
            slide = prs.slides.add_slide(slide_layout)
            if slide.shapes.title and s.get("title"):
                slide.shapes.title.text = s["title"]
            if len(slide.placeholders) > 1 and s.get("content"):
                slide.placeholders[1].text = s["content"]

        prs.save(output_path)
        return json.dumps({"status": "success", "file": output_path, "slides_count": len(prs.slides)})
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
def read_presentation(file_path: str) -> str:
    """Read a PowerPoint file and extract text from all slides.

    Args:
        file_path: Path to .pptx file.
    """
    if not PPTX_AVAILABLE:
        return json.dumps({"error": "python-pptx not installed"})
    try:
        prs = Presentation(file_path)
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
            slides.append({"slide_number": i + 1, "text": "\n".join(texts)})
        return json.dumps({"slides": slides, "total_slides": len(slides)}, indent=2)
    except Exception as e:
        return json.dumps({"error": str(e)})


@mcp.tool()
def get_presentation_info(file_path: str) -> str:
    """Get metadata and structure info about a PowerPoint file.

    Args:
        file_path: Path to .pptx file.
    """
    if not PPTX_AVAILABLE:
        return json.dumps({"error": "python-pptx not installed"})
    try:
        prs = Presentation(file_path)
        slide_info = []
        for i, slide in enumerate(prs.slides):
            shapes = [{"name": s.name, "type": s.shape_type.__class__.__name__ if hasattr(s.shape_type, '__class__') else str(s.shape_type)} for s in slide.shapes]
            slide_info.append({"slide_number": i + 1, "shapes_count": len(shapes), "shapes": shapes})
        return json.dumps({
            "total_slides": len(prs.slides),
            "slide_width": str(prs.slide_width),
            "slide_height": str(prs.slide_height),
            "slides": slide_info,
        }, indent=2, default=str)
    except Exception as e:
        return json.dumps({"error": str(e)})


if __name__ == "__main__":
    logger.info("Starting powerpoint-tools-mcp (transport=%s, port=%s)", MCP_TRANSPORT, MCP_PORT)
    if MCP_TRANSPORT == "stdio":
        mcp.run(transport="stdio", show_banner=False)
    else:
        mcp.run(transport="streamable-http", host="0.0.0.0", port=MCP_PORT)
