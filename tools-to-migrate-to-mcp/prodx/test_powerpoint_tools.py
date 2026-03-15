"""
Comprehensive tests for PowerPoint Tools

Tests all 3 PowerPoint tools:
- CreatePresentationTool
- AddSlideTool
- AddChartToSlideTool
"""

import pytest
import base64
import io
import json
import tempfile
import os

# Conditional imports for optional dependencies
try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    pytest.skip("python-pptx not available", allow_module_level=True)

# Import tools - handle both relative and absolute imports
try:
    from .powerpoint_tools import (
        CreatePresentationTool,
        AddSlideTool,
        AddChartToSlideTool,
        _decode_file_input,
        _encode_file_output
    )
except ImportError:
    # Fallback for direct execution
    import sys
    import os
    sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(__file__)))))
    from shared.modules.tools.prodx.powerpoint_tools import (
        CreatePresentationTool,
        AddSlideTool,
        AddChartToSlideTool,
        _decode_file_input,
        _encode_file_output
    )


class TestHelperFunctions:
    """Test helper functions."""
    
    def test_decode_file_input_base64(self):
        """Test decoding base64 file input."""
        test_data = b"test pptx content"
        encoded = base64.b64encode(test_data).decode('utf-8')
        
        file_bytes, is_base64 = _decode_file_input(encoded)
        assert file_bytes == test_data
        assert is_base64 is True
    
    def test_encode_file_output(self):
        """Test encoding file bytes to base64."""
        test_data = b"test content"
        encoded = _encode_file_output(test_data)
        decoded = base64.b64decode(encoded)
        assert decoded == test_data


class TestCreatePresentationTool:
    """Test CreatePresentationTool."""
    
    def test_create_presentation_basic(self):
        """Test creating a basic presentation."""
        tool = CreatePresentationTool()
        
        slides = [
            {
                "layout": "content",
                "title": "Slide 1",
                "content": ["Point 1", "Point 2"]
            },
            {
                "layout": "content",
                "title": "Slide 2",
                "content": ["Point 3", "Point 4"]
            }
        ]
        
        result = tool._run("Test Presentation", slides)
        result_dict = json.loads(result)
        
        assert result_dict["status"] == "success"
        assert "file_data" in result_dict
        assert result_dict["file_name"].endswith(".pptx")
    
    def test_create_presentation_with_images(self):
        """Test creating presentation with images."""
        tool = CreatePresentationTool()
        
        # Create a simple test image (1x1 pixel PNG)
        test_image = b'\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\tpHYs\x00\x00\x0b\x13\x00\x00\x0b\x13\x01\x00\x9a\x9c\x18\x00\x00\x00\nIDATx\x9cc\xf8\x00\x00\x00\x01\x00\x01\x00\x00\x00\x00IEND\xaeB`\x82'
        image_b64 = base64.b64encode(test_image).decode('utf-8')
        
        slides = [
            {
                "layout": "content",
                "title": "Slide with Image",
                "content": ["Test"],
                "images": [image_b64]
            }
        ]
        
        result = tool._run("Test Presentation", slides)
        result_dict = json.loads(result)
        assert result_dict["status"] == "success"
    
    def test_create_presentation_different_layouts(self):
        """Test creating presentation with different layouts."""
        tool = CreatePresentationTool()
        
        slides = [
            {"layout": "title", "title": "Title Slide"},
            {"layout": "content", "title": "Content Slide", "content": ["Text"]},
            {"layout": "two_content", "title": "Two Content", "content": ["Left", "Right"]},
            {"layout": "blank", "title": "Blank"}
        ]
        
        result = tool._run("Test", slides)
        result_dict = json.loads(result)
        assert result_dict["status"] == "success"
    
    def test_create_presentation_empty_slides(self):
        """Test creating presentation with empty slides list."""
        tool = CreatePresentationTool()
        
        result = tool._run("Test", [])
        result_dict = json.loads(result)
        # Should still create title slide
        assert result_dict["status"] == "success"


class TestAddSlideTool:
    """Test AddSlideTool."""
    
    def create_test_presentation(self):
        """Create a test PowerPoint presentation."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Existing Slide"
        
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return base64.b64encode(output.read()).decode('utf-8')
    
    def test_add_slide_basic(self):
        """Test adding a basic slide."""
        tool = AddSlideTool()
        pptx_data = self.create_test_presentation()
        
        slide_config = {
            "layout": "content",
            "title": "New Slide",
            "content": ["Point 1", "Point 2"]
        }
        
        result = tool._run(pptx_data, slide_config)
        result_dict = json.loads(result)
        
        assert result_dict["status"] == "success"
        assert "file_data" in result_dict
    
    def test_add_slide_at_position(self):
        """Test adding slide at specific position."""
        tool = AddSlideTool()
        pptx_data = self.create_test_presentation()
        
        slide_config = {
            "layout": "content",
            "title": "Inserted Slide",
            "content": ["Test"]
        }
        
        result = tool._run(pptx_data, slide_config, position=0)
        result_dict = json.loads(result)
        assert result_dict["status"] == "success"
    
    def test_add_slide_invalid_position(self):
        """Test adding slide at invalid position."""
        tool = AddSlideTool()
        pptx_data = self.create_test_presentation()
        
        slide_config = {
            "layout": "content",
            "title": "Test",
            "content": ["Test"]
        }
        
        # Position out of range should append
        result = tool._run(pptx_data, slide_config, position=999)
        result_dict = json.loads(result)
        assert result_dict["status"] == "success"
    
    def test_add_slide_with_images(self):
        """Test adding slide with images."""
        tool = AddSlideTool()
        pptx_data = self.create_test_presentation()
        
        test_image = b'\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\tpHYs\x00\x00\x0b\x13\x00\x00\x0b\x13\x01\x00\x9a\x9c\x18\x00\x00\x00\nIDATx\x9cc\xf8\x00\x00\x00\x01\x00\x01\x00\x00\x00\x00IEND\xaeB`\x82'
        image_b64 = base64.b64encode(test_image).decode('utf-8')
        
        slide_config = {
            "layout": "content",
            "title": "Slide with Image",
            "content": ["Test"],
            "images": [image_b64]
        }
        
        result = tool._run(pptx_data, slide_config)
        result_dict = json.loads(result)
        assert result_dict["status"] == "success"


class TestAddChartToSlideTool:
    """Test AddChartToSlideTool."""
    
    def create_test_presentation(self):
        """Create a test PowerPoint presentation."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Chart Slide"
        
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return base64.b64encode(output.read()).decode('utf-8')
    
    def test_add_chart_basic(self):
        """Test adding a chart to slide."""
        tool = AddChartToSlideTool()
        pptx_data = self.create_test_presentation()
        
        chart_config = {
            "chart_type": "bar",
            "data": [
                {"x": "A", "y": 10},
                {"x": "B", "y": 20},
                {"x": "C", "y": 30}
            ],
            "title": "Test Chart"
        }
        
        result = tool._run(pptx_data, 0, chart_config)
        result_dict = json.loads(result)
        
        assert result_dict["status"] == "success"
        assert "note" in result_dict  # Placeholder note about matplotlib/plotly
    
    def test_add_chart_invalid_slide_index(self):
        """Test adding chart to invalid slide index."""
        tool = AddChartToSlideTool()
        pptx_data = self.create_test_presentation()
        
        chart_config = {
            "chart_type": "bar",
            "data": [{"x": "A", "y": 10}]
        }
        
        # Invalid slide index
        result = tool._run(pptx_data, 999, chart_config)
        assert "Error" in result
    
    def test_add_chart_negative_index(self):
        """Test adding chart with negative slide index."""
        tool = AddChartToSlideTool()
        pptx_data = self.create_test_presentation()
        
        chart_config = {
            "chart_type": "bar",
            "data": [{"x": "A", "y": 10}]
        }
        
        result = tool._run(pptx_data, -1, chart_config)
        assert "Error" in result


class TestRealWorldScenarios:
    """Real-world scenario tests for PowerPoint tools."""
    
    def test_create_presentation_business_report(self):
        """Test creating a realistic business presentation."""
        tool = CreatePresentationTool()
        
        slides = [
            {
                "layout": "title",
                "title": "Q4 2024 Business Report"
            },
            {
                "layout": "content",
                "title": "Executive Summary",
                "content": [
                    "Revenue increased by 15%",
                    "Expenses decreased by 5%",
                    "Net profit margin improved to 25%"
                ]
            },
            {
                "layout": "two_content",
                "title": "Key Metrics",
                "content": [
                    {"text": "Revenue: $2.5M", "level": 0},
                    {"text": "Expenses: $1.5M", "level": 0},
                    {"text": "Profit: $1.0M", "level": 0}
                ]
            },
            {
                "layout": "content",
                "title": "Next Steps",
                "content": [
                    "Expand to new markets",
                    "Invest in R&D",
                    "Strengthen partnerships"
                ]
            }
        ]
        
        result = tool._run("Q4 2024 Business Report", slides)
        result_dict = json.loads(result)
        assert result_dict["status"] == "success"
        assert result_dict.get("slides_added", 0) >= 4
    
    def test_add_multiple_slides_sequentially(self):
        """Test adding multiple slides to existing presentation."""
        tool = CreatePresentationTool()
        
        # Create base presentation
        base_slides = [
            {
                "layout": "title",
                "title": "Base Presentation"
            }
        ]
        
        result = tool._run("Base", base_slides)
        result_dict = json.loads(result)
        assert result_dict["status"] == "success"
        base_pptx = json.loads(result_dict["file_data"]) if "file_data" in result_dict else result_dict.get("file_data", "")
        
        if base_pptx:
            # Add multiple slides
            add_tool = AddSlideTool()
            for i in range(5):
                slide_config = {
                    "layout": "content",
                    "title": f"Slide {i + 2}",
                    "content": [f"Content for slide {i + 2}"]
                }
                result = add_tool._run(base_pptx, slide_config)
                result_dict = json.loads(result)
                assert result_dict["status"] == "success"
                base_pptx = json.loads(result_dict["file_data"]) if "file_data" in result_dict else result_dict.get("file_data", "")
    
    def test_presentation_with_mixed_content(self):
        """Test presentation with various content types."""
        tool = CreatePresentationTool()
        
        # Create test image
        from PIL import Image
        img = Image.new('RGB', (400, 300), color='blue')
        buffer = io.BytesIO()
        img.save(buffer, format='PNG')
        buffer.seek(0)
        image_b64 = base64.b64encode(buffer.read()).decode('utf-8')
        
        slides = [
            {
                "layout": "content",
                "title": "Text Only Slide",
                "content": ["Point 1", "Point 2", "Point 3"]
            },
            {
                "layout": "content",
                "title": "Slide with Image",
                "content": ["This slide has an image"],
                "images": [image_b64]
            },
            {
                "layout": "blank",
                "title": "Blank Layout"
            }
        ]
        
        result = tool._run("Mixed Content", slides)
        result_dict = json.loads(result)
        assert result_dict["status"] == "success"


class TestErrorHandling:
    """Test error handling across all tools."""
    
    def test_all_tools_handle_invalid_base64(self):
        """Test that all tools handle invalid base64."""
        tools = [
            CreatePresentationTool(),
            AddSlideTool(),
            AddChartToSlideTool()
        ]
        
        for tool in tools:
            if isinstance(tool, CreatePresentationTool):
                result = tool._run("Test", [])
            elif isinstance(tool, AddSlideTool):
                result = tool._run("invalid_base64", {"layout": "content", "title": "Test"})
            else:  # AddChartToSlideTool
                result = tool._run("invalid_base64", 0, {"chart_type": "bar", "data": []})
            
            assert isinstance(result, str)
            assert len(result) > 0
    
    def test_all_tools_handle_empty_input(self):
        """Test that all tools handle empty input."""
        tool = CreatePresentationTool()
        result = tool._run("", [])
        assert isinstance(result, str)
        
        tool = AddSlideTool()
        result = tool._run("", {"layout": "content", "title": ""})
        assert isinstance(result, str)
    
    def test_all_tools_handle_corrupted_pptx(self):
        """Test that tools handle corrupted PowerPoint files."""
        tools = [AddSlideTool(), AddChartToSlideTool()]
        
        # Create corrupted base64 (valid base64 but not valid PPTX)
        corrupted_data = base64.b64encode(b"not a pptx file").decode('utf-8')
        
        for tool in tools:
            if isinstance(tool, AddSlideTool):
                result = tool._run(corrupted_data, {"layout": "content", "title": "Test"})
            else:
                result = tool._run(corrupted_data, 0, {"chart_type": "bar", "data": []})
            
            assert isinstance(result, str)
            assert "Error" in result


if __name__ == "__main__":
    pytest.main([__file__, "-v", "--tb=short"])

