"""
Comprehensive tests for File Operations Tools

Tests both file operations tools:
- SaveFileForDownloadTool
- ConvertFileFormatTool
"""

import pytest
import base64
import io
import json
import tempfile
import os

# Conditional imports for optional dependencies
try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False
    pytest.skip("pandas not available", allow_module_level=True)

try:
    from openpyxl import Workbook
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Import tools - handle both relative and absolute imports
try:
    from .file_operations_tools import (
        SaveFileForDownloadTool,
        ConvertFileFormatTool,
        _decode_file_input,
        _encode_file_output
    )
except ImportError:
    # Fallback for direct execution
    import sys
    import os
    sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(__file__)))))
    from shared.modules.tools.prodx.file_operations_tools import (
        SaveFileForDownloadTool,
        ConvertFileFormatTool,
        _decode_file_input,
        _encode_file_output
    )


class TestHelperFunctions:
    """Test helper functions."""
    
    def test_decode_file_input_base64(self):
        """Test decoding base64 file input."""
        test_data = b"test file content"
        encoded = base64.b64encode(test_data).decode('utf-8')
        
        file_bytes, is_base64 = _decode_file_input(encoded)
        assert file_bytes == test_data
        assert is_base64 is True
    
    def test_decode_file_input_file_path(self):
        """Test decoding file path input."""
        with tempfile.NamedTemporaryFile(delete=False) as tmp:
            tmp.write(b"test content")
            tmp_path = tmp.name
        
        try:
            file_bytes, is_base64 = _decode_file_input(tmp_path)
            assert file_bytes == b"test content"
            assert is_base64 is False
        finally:
            os.unlink(tmp_path)
    
    def test_encode_file_output(self):
        """Test encoding file bytes to base64."""
        test_data = b"test content"
        encoded = _encode_file_output(test_data)
        decoded = base64.b64decode(encoded)
        assert decoded == test_data


class TestSaveFileForDownloadTool:
    """Test SaveFileForDownloadTool."""
    
    def test_save_file_basic(self):
        """Test basic file saving."""
        tool = SaveFileForDownloadTool()
        
        test_data = b"test file content"
        file_data = base64.b64encode(test_data).decode('utf-8')
        
        result = tool._run(
            file_data,
            "test_file.txt",
            "text/plain"
        )
        result_dict = json.loads(result)
        
        assert result_dict["status"] == "success"
        assert "file_id" in result_dict or "message" in result_dict
    
    def test_save_file_excel(self):
        """Test saving Excel file."""
        if not OPENPYXL_AVAILABLE:
            pytest.skip("openpyxl not available")
        
        tool = SaveFileForDownloadTool()
        
        # Create a simple Excel file
        wb = Workbook()
        ws = wb.active
        ws['A1'] = "Test"
        
        output = io.BytesIO()
        wb.save(output)
        output.seek(0)
        file_data = base64.b64encode(output.read()).decode('utf-8')
        
        result = tool._run(
            file_data,
            "test.xlsx",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )
        result_dict = json.loads(result)
        assert result_dict["status"] == "success"
    
    def test_save_file_powerpoint(self):
        """Test saving PowerPoint file."""
        if not PPTX_AVAILABLE:
            pytest.skip("python-pptx not available")
        
        tool = SaveFileForDownloadTool()
        
        # Create a simple PPTX file
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Test"
        
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        file_data = base64.b64encode(output.read()).decode('utf-8')
        
        result = tool._run(
            file_data,
            "test.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        result_dict = json.loads(result)
        assert result_dict["status"] == "success"
    
    def test_save_file_invalid_base64(self):
        """Test saving with invalid base64."""
        tool = SaveFileForDownloadTool()
        
        result = tool._run("invalid_base64", "test.txt", "text/plain")
        assert "Error" in result
    
    def test_save_file_with_session_state(self):
        """Test saving file to session state."""
        tool = SaveFileForDownloadTool()
        
        test_data = b"test content"
        file_data = base64.b64encode(test_data).decode('utf-8')
        
        result = tool._run(
            file_data,
            "test.txt",
            "text/plain",
            storage_location="session_state"
        )
        result_dict = json.loads(result)
        assert result_dict["status"] == "success"


class TestConvertFileFormatTool:
    """Test ConvertFileFormatTool."""
    
    def create_test_csv(self):
        """Create a test CSV file."""
        df = pd.DataFrame({
            "Name": ["Alice", "Bob"],
            "Age": [30, 25]
        })
        
        output = io.BytesIO()
        df.to_csv(output, index=False)
        output.seek(0)
        return base64.b64encode(output.read()).decode('utf-8')
    
    def create_test_json(self):
        """Create a test JSON file."""
        data = {
            "data": [
                {"name": "Alice", "age": 30},
                {"name": "Bob", "age": 25}
            ]
        }
        json_str = json.dumps(data)
        return base64.b64encode(json_str.encode('utf-8')).decode('utf-8')
    
    def test_convert_csv_to_json(self):
        """Test converting CSV to JSON."""
        tool = ConvertFileFormatTool()
        csv_data = self.create_test_csv()
        
        result = tool._run(csv_data, "csv", "json")
        result_dict = json.loads(result)
        
        assert result_dict["status"] == "success"
        assert "file_data" in result_dict
        assert result_dict["file_name"].endswith(".json")
    
    def test_convert_json_to_csv(self):
        """Test converting JSON to CSV."""
        tool = ConvertFileFormatTool()
        json_data = self.create_test_json()
        
        result = tool._run(json_data, "json", "csv")
        result_dict = json.loads(result)
        
        assert result_dict["status"] == "success"
        assert "file_data" in result_dict
        assert result_dict["file_name"].endswith(".csv")
    
    def test_convert_excel_to_csv(self):
        """Test converting Excel to CSV."""
        if not OPENPYXL_AVAILABLE:
            pytest.skip("openpyxl not available")
        
        tool = ConvertFileFormatTool()
        
        # Create Excel file
        wb = Workbook()
        ws = wb.active
        ws.append(["Name", "Age"])
        ws.append(["Alice", 30])
        ws.append(["Bob", 25])
        
        output = io.BytesIO()
        wb.save(output)
        output.seek(0)
        excel_data = base64.b64encode(output.read()).decode('utf-8')
        
        result = tool._run(excel_data, "excel", "csv")
        result_dict = json.loads(result)
        assert result_dict["status"] == "success"
    
    def test_convert_csv_to_excel(self):
        """Test converting CSV to Excel."""
        tool = ConvertFileFormatTool()
        csv_data = self.create_test_csv()
        
        result = tool._run(csv_data, "csv", "excel")
        result_dict = json.loads(result)
        assert result_dict["status"] == "success"
    
    def test_convert_unsupported_format(self):
        """Test converting unsupported format."""
        tool = ConvertFileFormatTool()
        test_data = base64.b64encode(b"test").decode('utf-8')
        
        result = tool._run(test_data, "unknown", "unknown")
        assert "Error" in result
    
    def test_convert_invalid_base64(self):
        """Test converting with invalid base64."""
        tool = ConvertFileFormatTool()
        
        result = tool._run("invalid_base64", "csv", "json")
        assert "Error" in result
    
    def test_convert_same_format(self):
        """Test converting same format (should handle gracefully)."""
        tool = ConvertFileFormatTool()
        csv_data = self.create_test_csv()
        
        result = tool._run(csv_data, "csv", "csv")
        # Should either succeed (no-op) or return appropriate message
        assert isinstance(result, str)


class TestErrorHandling:
    """Test error handling across all file operations tools."""
    
    def test_all_tools_handle_invalid_base64(self):
        """Test that all tools handle invalid base64."""
        tools = [
            SaveFileForDownloadTool(),
            ConvertFileFormatTool()
        ]
        
        for tool in tools:
            if isinstance(tool, SaveFileForDownloadTool):
                result = tool._run("invalid_base64", "test.txt", "text/plain")
            else:  # ConvertFileFormatTool
                result = tool._run("invalid_base64", "csv", "json")
            
            assert isinstance(result, str)
            assert len(result) > 0
    
    def test_all_tools_handle_empty_input(self):
        """Test that all tools handle empty input."""
        tools = [
            SaveFileForDownloadTool(),
            ConvertFileFormatTool()
        ]
        
        for tool in tools:
            if isinstance(tool, SaveFileForDownloadTool):
                result = tool._run("", "test.txt", "text/plain")
            else:  # ConvertFileFormatTool
                result = tool._run("", "csv", "json")
            
            assert isinstance(result, str)
    
    def test_all_tools_handle_missing_streamlit(self):
        """Test that tools handle missing streamlit gracefully."""
        # Tools should return error messages, not crash
        tool = SaveFileForDownloadTool()
        test_data = base64.b64encode(b"test").decode('utf-8')
        
        result = tool._run(test_data, "test.txt", "text/plain")
        assert isinstance(result, str)


if __name__ == "__main__":
    pytest.main([__file__, "-v", "--tb=short"])

