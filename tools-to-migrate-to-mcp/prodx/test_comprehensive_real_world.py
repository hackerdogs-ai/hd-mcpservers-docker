"""
Comprehensive Real-World Integration Tests

This file contains end-to-end tests that simulate real-world usage scenarios
with actual data, web downloads, and complex workflows.
"""

import pytest
import base64
import io
import json
import tempfile
import os
import requests
from typing import Optional

# Try to import all required libraries
try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

try:
    from openpyxl import Workbook
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

try:
    import requests
    REQUESTS_AVAILABLE = True
except ImportError:
    REQUESTS_AVAILABLE = False

# Import all tools
try:
    from .excel_tools import (
        ReadExcelStructuredTool,
        ModifyExcelTool,
        CreateExcelChartTool,
        AnalyzeExcelSecurityTool
    )
    from .powerpoint_tools import (
        CreatePresentationTool,
        AddSlideTool,
        AddChartToSlideTool
    )
    from .visualization_tools import (
        CreatePlotlyChartTool,
        CreateChartFromFileTool
    )
    from .ocr_tools import (
        ExtractTextFromImageTool,
        ExtractTextFromPDFImagesTool,
        AnalyzeDocumentStructureTool
    )
    from .file_operations_tools import (
        SaveFileForDownloadTool,
        ConvertFileFormatTool
    )
except ImportError:
    import sys
    import os
    sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(__file__)))))
    from shared.modules.tools.prodx.excel_tools import (
        ReadExcelStructuredTool,
        ModifyExcelTool,
        CreateExcelChartTool,
        AnalyzeExcelSecurityTool
    )
    from shared.modules.tools.prodx.powerpoint_tools import (
        CreatePresentationTool,
        AddSlideTool,
        AddChartToSlideTool
    )
    from shared.modules.tools.prodx.visualization_tools import (
        CreatePlotlyChartTool,
        CreateChartFromFileTool
    )
    from shared.modules.tools.prodx.ocr_tools import (
        ExtractTextFromImageTool,
        ExtractTextFromPDFImagesTool,
        AnalyzeDocumentStructureTool
    )
    from shared.modules.tools.prodx.file_operations_tools import (
        SaveFileForDownloadTool,
        ConvertFileFormatTool
    )


class TestEndToEndWorkflows:
    """End-to-end workflow tests simulating real user scenarios."""
    
    @pytest.mark.skipif(not OPENPYXL_AVAILABLE, reason="openpyxl not available")
    def test_excel_analysis_to_powerpoint_workflow(self):
        """Test complete workflow: Read Excel -> Analyze -> Create PowerPoint."""
        # Step 1: Create and read Excel file
        wb = Workbook()
        ws = wb.active
        ws.append(["Quarter", "Sales", "Expenses"])
        ws.append(["Q1", 100000, 60000])
        ws.append(["Q2", 120000, 70000])
        ws.append(["Q3", 110000, 65000])
        ws.append(["Q4", 130000, 75000])
        
        excel_output = io.BytesIO()
        wb.save(excel_output)
        excel_output.seek(0)
        excel_data = base64.b64encode(excel_output.read()).decode('utf-8')
        
        # Step 2: Read Excel
        read_tool = ReadExcelStructuredTool()
        read_result = read_tool._run(excel_data)
        read_dict = json.loads(read_result)
        assert read_dict["status"] == "success"
        
        # Step 3: Analyze security
        security_tool = AnalyzeExcelSecurityTool()
        security_result = security_tool._run(excel_data)
        security_dict = json.loads(security_result)
        assert security_dict["status"] == "success"
        
        # Step 4: Create chart in Excel
        chart_tool = CreateExcelChartTool()
        chart_config = {
            "chart_type": "bar",
            "data_range": "A1:C5",
            "title": "Sales Analysis"
        }
        chart_result = chart_tool._run(excel_data, chart_config)
        chart_dict = json.loads(chart_result)
        assert chart_dict["status"] == "success"
        
        # Step 5: Create PowerPoint with results
        if PPTX_AVAILABLE:
            ppt_tool = CreatePresentationTool()
            slides = [
                {
                    "layout": "content",
                    "title": "Sales Analysis Report",
                    "content": [
                        "Q1 Sales: $100,000",
                        "Q2 Sales: $120,000",
                        "Q3 Sales: $110,000",
                        "Q4 Sales: $130,000"
                    ]
                }
            ]
            ppt_result = ppt_tool._run("Sales Analysis", slides)
            ppt_dict = json.loads(ppt_result)
            assert ppt_dict["status"] == "success"
    
    @pytest.mark.skipif(not REQUESTS_AVAILABLE or not PIL_AVAILABLE, reason="requests or PIL not available")
    def test_image_ocr_to_excel_workflow(self):
        """Test workflow: Download image -> OCR -> Create Excel with results."""
        # Step 1: Download test image
        test_url = "https://via.placeholder.com/400x200/000000/FFFFFF?text=Test+Data+123"
        try:
            response = requests.get(test_url, timeout=10)
            response.raise_for_status()
            img = Image.open(io.BytesIO(response.content))
        except Exception as e:
            pytest.skip(f"Could not download test image: {e}")
        
        # Step 2: Extract text with OCR
        ocr_tool = ExtractTextFromImageTool()
        buffer = io.BytesIO()
        img.save(buffer, format='PNG')
        buffer.seek(0)
        image_data = base64.b64encode(buffer.read()).decode('utf-8')
        
        ocr_result = ocr_tool._run(image_data)
        ocr_dict = json.loads(ocr_result)
        assert ocr_dict["status"] == "success"
        
        # Step 3: Create Excel with OCR results
        if OPENPYXL_AVAILABLE:
            wb = Workbook()
            ws = wb.active
            ws.append(["Source", "Extracted Text"])
            ws.append(["Image OCR", ocr_dict.get("text", "")])
            
            excel_output = io.BytesIO()
            wb.save(excel_output)
            excel_output.seek(0)
            excel_data = base64.b64encode(excel_output.read()).decode('utf-8')
            
            # Verify Excel was created
            read_tool = ReadExcelStructuredTool()
            read_result = read_tool._run(excel_data)
            read_dict = json.loads(read_result)
            assert read_dict["status"] == "success"
    
    @pytest.mark.skipif(not PANDAS_AVAILABLE, reason="pandas not available")
    def test_csv_to_chart_to_powerpoint_workflow(self):
        """Test workflow: CSV -> Chart -> PowerPoint."""
        # Step 1: Create CSV data
        df = pd.DataFrame({
            "Month": ["Jan", "Feb", "Mar", "Apr"],
            "Revenue": [50000, 55000, 60000, 65000]
        })
        
        csv_output = io.BytesIO()
        df.to_csv(csv_output, index=False)
        csv_output.seek(0)
        csv_data = base64.b64encode(csv_output.read()).decode('utf-8')
        
        # Step 2: Create chart from CSV
        chart_tool = CreateChartFromFileTool()
        chart_result = chart_tool._run(csv_data, "csv", "line_chart", display=False)
        chart_dict = json.loads(chart_result)
        assert chart_dict["status"] == "success"
        
        # Step 3: Add chart slide to PowerPoint
        if PPTX_AVAILABLE:
            # Create base presentation
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Revenue Chart"
            
            ppt_output = io.BytesIO()
            prs.save(ppt_output)
            ppt_output.seek(0)
            ppt_data = base64.b64encode(ppt_output.read()).decode('utf-8')
            
            # Add chart slide
            add_slide_tool = AddSlideTool()
            slide_config = {
                "layout": "content",
                "title": "Revenue Analysis",
                "content": ["Revenue trend over 4 months"]
            }
            add_result = add_slide_tool._run(ppt_data, slide_config)
            add_dict = json.loads(add_result)
            assert add_dict["status"] == "success"


class TestPerformanceAndStress:
    """Performance and stress tests."""
    
    @pytest.mark.skipif(not OPENPYXL_AVAILABLE, reason="openpyxl not available")
    def test_large_excel_file_handling(self):
        """Test handling of large Excel files."""
        # Create large Excel file (1000 rows)
        wb = Workbook()
        ws = wb.active
        ws.append(["ID", "Name", "Value", "Category"])
        
        for i in range(1000):
            ws.append([i + 1, f"Item_{i+1}", (i + 1) * 10, f"Cat_{i % 10}"])
        
        excel_output = io.BytesIO()
        wb.save(excel_output)
        excel_output.seek(0)
        excel_data = base64.b64encode(excel_output.read()).decode('utf-8')
        
        # Test reading
        read_tool = ReadExcelStructuredTool()
        read_result = read_tool._run(excel_data)
        read_dict = json.loads(read_result)
        assert read_dict["status"] == "success"
        
        # Test modification
        modify_tool = ModifyExcelTool()
        operations = [
            {
                "type": "update_cell",
                "sheet": "Sheet1",
                "cell": "A1001",
                "value": "Total"
            }
        ]
        modify_result = modify_tool._run(excel_data, operations)
        modify_dict = json.loads(modify_result)
        assert modify_dict["status"] == "success"
    
    @pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
    def test_large_powerpoint_handling(self):
        """Test handling of large PowerPoint files."""
        # Create presentation with many slides
        prs = Presentation()
        for i in range(50):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            if slide.shapes.title:
                slide.shapes.title.text = f"Slide {i + 1}"
        
        ppt_output = io.BytesIO()
        prs.save(ppt_output)
        ppt_output.seek(0)
        ppt_data = base64.b64encode(ppt_output.read()).decode('utf-8')
        
        # Test adding slide
        add_slide_tool = AddSlideTool()
        slide_config = {
            "layout": "content",
            "title": "New Slide",
            "content": ["Test"]
        }
        result = add_slide_tool._run(ppt_data, slide_config)
        result_dict = json.loads(result)
        assert result_dict["status"] == "success"


class TestErrorRecovery:
    """Test error recovery and resilience."""
    
    def test_all_tools_recover_from_errors(self):
        """Test that all tools recover gracefully from various errors."""
        tools_to_test = []
        
        if OPENPYXL_AVAILABLE:
            tools_to_test.extend([
                (ReadExcelStructuredTool(), "invalid_base64"),
                (ModifyExcelTool(), "invalid_base64", []),
                (CreateExcelChartTool(), "invalid_base64", {}),
                (AnalyzeExcelSecurityTool(), "invalid_base64")
            ])
        
        if PPTX_AVAILABLE:
            tools_to_test.extend([
                (CreatePresentationTool(), "Test", []),
                (AddSlideTool(), "invalid_base64", {}),
                (AddChartToSlideTool(), "invalid_base64", 0, {})
            ])
        
        tools_to_test.extend([
            (CreatePlotlyChartTool(), {}, "line_chart"),
            (CreateChartFromFileTool(), "invalid_base64", "csv", "line_chart")
        ])
        
        if PIL_AVAILABLE:
            tools_to_test.extend([
                (ExtractTextFromImageTool(), "invalid_base64"),
                (ExtractTextFromPDFImagesTool(), "invalid_base64"),
                (AnalyzeDocumentStructureTool(), "invalid_base64")
            ])
        
        tools_to_test.extend([
            (SaveFileForDownloadTool(), "invalid_base64", "test.txt", "text/plain"),
            (ConvertFileFormatTool(), "invalid_base64", "csv", "json")
        ])
        
        for tool_data in tools_to_test:
            tool = tool_data[0]
            args = tool_data[1:]
            
            try:
                result = tool._run(*args)
                # Should return error message, not crash
                assert isinstance(result, str)
                assert len(result) > 0
            except Exception as e:
                # Tools should catch exceptions internally
                pytest.fail(f"Tool {type(tool).__name__} raised exception: {e}")


class TestDataIntegrity:
    """Test data integrity and correctness."""
    
    @pytest.mark.skipif(not OPENPYXL_AVAILABLE, reason="openpyxl not available")
    def test_excel_data_preservation(self):
        """Test that Excel data is preserved through operations."""
        # Create Excel with specific data
        wb = Workbook()
        ws = wb.active
        ws['A1'] = "Original Value"
        ws['B1'] = 12345
        ws['A2'] = "Test"
        ws['B2'] = 67890
        
        excel_output = io.BytesIO()
        wb.save(excel_output)
        excel_output.seek(0)
        excel_data = base64.b64encode(excel_output.read()).decode('utf-8')
        
        # Read and verify
        read_tool = ReadExcelStructuredTool()
        read_result = read_tool._run(excel_data)
        read_dict = json.loads(read_result)
        assert read_dict["status"] == "success"
        
        # Modify and verify data is updated
        modify_tool = ModifyExcelTool()
        operations = [
            {
                "type": "update_cell",
                "sheet": "Sheet1",
                "cell": "A1",
                "value": "Updated Value"
            }
        ]
        modify_result = modify_tool._run(excel_data, operations)
        modify_dict = json.loads(modify_result)
        assert modify_dict["status"] == "success"
        
        # Read modified file
        modified_data = json.loads(modify_dict["file_data"]) if "file_data" in modify_dict else None
        if modified_data:
            # Decode and verify
            modified_bytes = base64.b64decode(modified_data)
            modified_wb = Workbook()
            modified_wb.load(io.BytesIO(modified_bytes))
            modified_ws = modified_wb.active
            assert modified_ws['A1'].value == "Updated Value"
            assert modified_ws['B1'].value == 12345  # Should be preserved


if __name__ == "__main__":
    pytest.main([__file__, "-v", "--tb=short"])

