"""
Utility functions for testing prodx tools.

Provides helper functions for creating test data, downloading test files,
and common test patterns.
"""

import base64
import io
import tempfile
import os
import requests
from typing import Optional, Tuple
from PIL import Image
import numpy as np

try:
    from openpyxl import Workbook
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False


def download_image_from_url(url: str, timeout: int = 10) -> Optional[Image.Image]:
    """
    Download an image from a URL for testing.
    
    Args:
        url: URL of the image to download
        timeout: Request timeout in seconds
        
    Returns:
        PIL Image or None if download fails
    """
    try:
        response = requests.get(url, timeout=timeout, stream=True)
        response.raise_for_status()
        
        img = Image.open(io.BytesIO(response.content))
        return img
    except Exception as e:
        print(f"Failed to download image from {url}: {e}")
        return None


def create_test_excel_file(data: Optional[list] = None, sheet_name: str = "Sheet1") -> Tuple[bytes, str]:
    """
    Create a test Excel file.
    
    Args:
        data: List of rows (each row is a list of values)
        sheet_name: Name of the sheet
        
    Returns:
        Tuple of (file_bytes, base64_encoded_string)
    """
    if not OPENPYXL_AVAILABLE:
        raise ImportError("openpyxl not available")
    
    wb = Workbook()
    ws = wb.active
    ws.title = sheet_name
    
    if data is None:
        data = [
            ["Name", "Age", "City", "Score"],
            ["Alice", 30, "New York", 95],
            ["Bob", 25, "London", 87],
            ["Charlie", 35, "Tokyo", 92],
            ["Diana", 28, "Paris", 88],
            ["Eve", 32, "Berlin", 90]
        ]
    
    for row in data:
        ws.append(row)
    
    output = io.BytesIO()
    wb.save(output)
    output.seek(0)
    file_bytes = output.read()
    encoded = base64.b64encode(file_bytes).decode('utf-8')
    
    return file_bytes, encoded


def create_test_powerpoint_file(slides: int = 1) -> Tuple[bytes, str]:
    """
    Create a test PowerPoint file.
    
    Args:
        slides: Number of slides to create
        
    Returns:
        Tuple of (file_bytes, base64_encoded_string)
    """
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx not available")
    
    prs = Presentation()
    
    for i in range(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        if slide.shapes.title:
            slide.shapes.title.text = f"Slide {i + 1}"
        if len(slide.placeholders) > 1:
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.text = f"Content for slide {i + 1}"
    
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    file_bytes = output.read()
    encoded = base64.b64encode(file_bytes).decode('utf-8')
    
    return file_bytes, encoded


def create_test_image_with_text(text: str = "TEST", size: Tuple[int, int] = (200, 100)) -> Tuple[bytes, str]:
    """
    Create a test image with text (simplified - actual text rendering requires PIL ImageDraw).
    
    Args:
        text: Text to include (not actually rendered, just for reference)
        size: Image size (width, height)
        
    Returns:
        Tuple of (file_bytes, base64_encoded_string)
    """
    img = Image.new('RGB', size, color='white')
    
    # Add some pattern to make it more realistic
    pixels = np.array(img)
    # Add a simple gradient
    for y in range(size[1]):
        for x in range(size[0]):
            pixels[y, x] = [min(255, x % 256), min(255, y % 256), 200]
    
    img = Image.fromarray(pixels)
    
    buffer = io.BytesIO()
    img.save(buffer, format='PNG')
    buffer.seek(0)
    file_bytes = buffer.read()
    encoded = base64.b64encode(file_bytes).decode('utf-8')
    
    return file_bytes, encoded


def create_test_csv_file(data: Optional[list] = None) -> Tuple[bytes, str]:
    """
    Create a test CSV file.
    
    Args:
        data: List of rows (each row is a list of values)
        
    Returns:
        Tuple of (file_bytes, base64_encoded_string)
    """
    if not PANDAS_AVAILABLE:
        raise ImportError("pandas not available")
    
    if data is None:
        data = {
            "Name": ["Alice", "Bob", "Charlie"],
            "Age": [30, 25, 35],
            "City": ["New York", "London", "Tokyo"]
        }
        df = pd.DataFrame(data)
    else:
        df = pd.DataFrame(data)
    
    output = io.BytesIO()
    df.to_csv(output, index=False)
    output.seek(0)
    file_bytes = output.read()
    encoded = base64.b64encode(file_bytes).decode('utf-8')
    
    return file_bytes, encoded


def create_test_json_file(data: Optional[dict] = None) -> Tuple[bytes, str]:
    """
    Create a test JSON file.
    
    Args:
        data: Dictionary to serialize
        
    Returns:
        Tuple of (file_bytes, base64_encoded_string)
    """
    import json
    
    if data is None:
        data = {
            "users": [
                {"name": "Alice", "age": 30, "city": "New York"},
                {"name": "Bob", "age": 25, "city": "London"},
                {"name": "Charlie", "age": 35, "city": "Tokyo"}
            ]
        }
    
    json_str = json.dumps(data, indent=2)
    file_bytes = json_str.encode('utf-8')
    encoded = base64.b64encode(file_bytes).decode('utf-8')
    
    return file_bytes, encoded


def get_test_image_urls() -> list:
    """
    Get URLs of test images for OCR testing.
    Returns a list of publicly available test images.
    """
    return [
        "https://via.placeholder.com/400x200/000000/FFFFFF?text=TEST+IMAGE",
        "https://httpbin.org/image/png",  # Returns a PNG image
        "https://httpbin.org/image/jpeg",  # Returns a JPEG image
    ]


def create_large_excel_file(rows: int = 1000) -> Tuple[bytes, str]:
    """
    Create a large Excel file for performance testing.
    
    Args:
        rows: Number of data rows
        
    Returns:
        Tuple of (file_bytes, base64_encoded_string)
    """
    if not OPENPYXL_AVAILABLE:
        raise ImportError("openpyxl not available")
    
    wb = Workbook()
    ws = wb.active
    ws.append(["ID", "Name", "Value", "Category", "Date"])
    
    for i in range(rows):
        ws.append([
            i + 1,
            f"Item_{i + 1}",
            (i + 1) * 10,
            f"Category_{i % 10}",
            f"2024-01-{(i % 28) + 1:02d}"
        ])
    
    output = io.BytesIO()
    wb.save(output)
    output.seek(0)
    file_bytes = output.read()
    encoded = base64.b64encode(file_bytes).decode('utf-8')
    
    return file_bytes, encoded


def create_excel_with_formulas() -> Tuple[bytes, str]:
    """
    Create an Excel file with formulas for testing.
    
    Returns:
        Tuple of (file_bytes, base64_encoded_string)
    """
    if not OPENPYXL_AVAILABLE:
        raise ImportError("openpyxl not available")
    
    wb = Workbook()
    ws = wb.active
    
    # Add data
    ws.append(["Item", "Price", "Quantity", "Total"])
    ws.append(["Apple", 1.50, 10, "=B2*C2"])
    ws.append(["Banana", 0.75, 20, "=B3*C3"])
    ws.append(["Orange", 2.00, 15, "=B4*C4"])
    ws.append(["", "", "SUM", "=SUM(D2:D4)"])
    
    output = io.BytesIO()
    wb.save(output)
    output.seek(0)
    file_bytes = output.read()
    encoded = base64.b64encode(file_bytes).decode('utf-8')
    
    return file_bytes, encoded


def create_excel_with_charts_data() -> Tuple[bytes, str]:
    """
    Create an Excel file with data suitable for chart creation.
    
    Returns:
        Tuple of (file_bytes, base64_encoded_string)
    """
    if not OPENPYXL_AVAILABLE:
        raise ImportError("openpyxl not available")
    
    wb = Workbook()
    ws = wb.active
    
    # Sales data by month
    ws.append(["Month", "Sales", "Expenses", "Profit"])
    months = ["Jan", "Feb", "Mar", "Apr", "May", "Jun"]
    for i, month in enumerate(months):
        sales = (i + 1) * 1000
        expenses = (i + 1) * 600
        profit = sales - expenses
        ws.append([month, sales, expenses, profit])
    
    output = io.BytesIO()
    wb.save(output)
    output.seek(0)
    file_bytes = output.read()
    encoded = base64.b64encode(file_bytes).decode('utf-8')
    
    return file_bytes, encoded

