"""
PowerPoint Operations Tools
---------------------------
LangChain tools for PowerPoint presentation generation and manipulation.
"""

from langchain_core.tools import BaseTool
from pydantic import BaseModel, Field
from typing import Type, Dict, List, Optional, Any
import sys
import os
# Add project root to path for shared modules
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '../../../../'))
try:
    from shared.logger import setup_logger
except ImportError:
    # Fallback for testing environments
    import logging
    def setup_logger(name, log_file_path=None, **kwargs):
        logger = logging.getLogger(name)
        if not logger.handlers:
            handler = logging.StreamHandler()
            handler.setFormatter(logging.Formatter('%(name)s - %(levelname)s - %(message)s'))
            logger.addHandler(handler)
            logger.setLevel(logging.INFO)
        return logger
import base64
import io
import json

logger = setup_logger(__name__, log_file_path="logs/powerpoint_tools.log")

# Try to import required libraries
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available. PowerPoint tools will not work.")


def _decode_file_input(file_input: str) -> tuple[bytes, bool]:
    """Decode file input which can be either a base64 string or file path."""
    logger.debug(f"Attempting to decode file input (length: {len(file_input) if file_input else 0})")
    
    try:
        # Try to decode as base64 first
        logger.debug("Attempting base64 decoding")
        file_bytes = base64.b64decode(file_input)
        logger.debug(f"Successfully decoded base64 input (size: {len(file_bytes)} bytes)")
        return file_bytes, True
    except Exception as base64_error:
        logger.debug(f"Base64 decoding failed: {str(base64_error)}, treating as file path")
        # If that fails, treat as file path
        try:
            logger.debug(f"Attempting to read file from path: {file_input[:100] if len(file_input) > 100 else file_input}")
            with open(file_input, 'rb') as f:
                file_bytes = f.read()
                logger.debug(f"Successfully read file from path (size: {len(file_bytes)} bytes)")
                return file_bytes, False
        except FileNotFoundError as fnf_error:
            logger.error(f"File not found: {file_input}")
            raise ValueError(f"File not found: {file_input}")
        except PermissionError as perm_error:
            logger.error(f"Permission denied reading file: {file_input}")
            raise ValueError(f"Permission denied reading file: {file_input}")
        except Exception as file_error:
            logger.error(f"Error reading file from path: {str(file_error)}", exc_info=True)
            raise ValueError(f"Invalid file input: {file_error}")


def _encode_file_output(file_bytes: bytes) -> str:
    """Encode file bytes to base64 string."""
    try:
        logger.debug(f"Encoding {len(file_bytes)} bytes to base64")
        encoded = base64.b64encode(file_bytes).decode('utf-8')
        logger.debug(f"Successfully encoded to base64 (length: {len(encoded)} chars)")
        return encoded
    except Exception as e:
        logger.error(f"Error encoding file bytes to base64: {str(e)}", exc_info=True)
        raise ValueError(f"Failed to encode file to base64: {str(e)}")


# --- PowerPoint Tools ---

class CreatePresentationInput(BaseModel):
    """Input schema for CreatePresentationTool."""
    title: str = Field(..., description="Presentation title")
    slides: List[Dict[str, Any]] = Field(..., description="List of slide configurations")
    template: Optional[str] = Field(None, description="Template name or path (optional)")
    output_format: str = Field("base64", description="Output format: 'base64' or 'file_path' (default: 'base64')")


class CreatePresentationTool(BaseTool):
    """
    Create PowerPoint presentations from structured data.
    
    This tool generates complete PowerPoint presentations from structured slide configurations.
    It supports multiple slide layouts (title, content, two-column, blank), text content with
    formatting, image embedding, and chart integration. The tool creates professional-looking
    presentations suitable for reports, analysis summaries, and executive briefings.
    
    **When to use this tool:**
    - Generating presentation decks from analysis results
    - Creating security analysis presentations
    - Generating executive summary reports in PowerPoint format
    - Creating structured presentations from data
    - Building report decks with multiple slides
    - Converting analysis results into presentation format
    - Creating slide decks for meetings or reviews
    
    **When NOT to use this tool:**
    - For modifying existing presentations (use AddSlideTool)
    - When you need complex animations or transitions (use PowerPoint directly)
    - For presentations requiring custom templates (template support is limited)
    - When you only need a single slide (consider AddSlideTool instead)
    
    **Input requirements:**
    - Must provide a presentation title
    - Slides list must contain valid slide configuration dictionaries
    - Each slide config should include layout type and content
    - Optional template path (if not provided, uses default template)
    - Images must be base64-encoded strings
    
    **Output:**
    - Returns JSON with status, base64-encoded PowerPoint file, file name, and MIME type
    - Presentation is ready for download or further processing
    - Returns error message if creation fails
    
    **Supported Slide Layouts:**
    - title: Title slide with title and subtitle
    - content: Title and content layout (most common)
    - two_content: Two-column content layout
    - blank: Blank slide for custom content
    
    **Content Format:**
    - Simple text: List of strings, each becomes a paragraph
    - Structured: List of dicts with "text" and optional "level" for bullet hierarchy
    
    **Limitations:**
    - Chart embedding requires matplotlib/plotly integration (placeholder)
    - Complex formatting options are limited
    - Animations and transitions not supported
    - Template support is basic (uses default PowerPoint templates)
    
    **Example use cases:**
    1. "Create a security analysis presentation with findings and recommendations"
    2. "Generate a PowerPoint deck from the threat analysis results"
    3. "Create an executive summary presentation with 5 slides"
    4. "Build a presentation deck with title slide and 3 content slides"
    
    **Configuration:**
    Requires python-pptx library to be installed. The tool automatically creates a title slide
    and adds content slides based on the provided configuration. Output is base64-encoded for
    easy integration with Streamlit downloads.
    """
    name: str = "create_presentation_deck"
    description: str = (
        "Create PowerPoint presentations from structured data. "
        "Supports multiple slide layouts (title, content, two_content, blank), text content, images, and charts. "
        "Returns base64-encoded PowerPoint file ready for download. "
        "IMPORTANT: Title and slides list are required. Images must be base64-encoded. "
        "Best for: generating presentation decks from analysis, creating reports, executive summaries. "
        "NOT suitable for: modifying existing presentations, complex animations, or custom template requirements."
    )
    args_schema: Type[BaseModel] = CreatePresentationInput

    def _run(
        self,
        title: str,
        slides: List[Dict[str, Any]],
        template: Optional[str] = None,
        output_format: str = "base64"
    ) -> str:
        """
        Create PowerPoint presentation.
        
        This method performs the presentation creation:
        1. Creates new Presentation object (or loads template if provided)
        2. Adds title slide with presentation title
        3. Iterates through slides list, creating each slide
        4. Applies layout, sets title, adds content
        5. Embeds images if provided
        6. Saves presentation to bytes
        7. Returns base64-encoded file or saves to path
        
        Args:
            title: Presentation title (appears on title slide).
                  Example: "Security Analysis Report", "Q4 Results"
            slides: List of slide configuration dictionaries. Each slide config should have:
                  - layout: Layout type (title, content, two_content, blank)
                  - title: Slide title
                  - content: List of content blocks (strings or dicts with text/level)
                  - images: Optional list of base64-encoded images
                  Example: [{"layout": "content", "title": "Findings", "content": ["Finding 1", "Finding 2"]}]
            template: Optional template file path (.pptx file)
            output_format: Output format - "base64" (default) or "file_path"
        
        Returns:
            str: JSON string with status, base64-encoded PowerPoint file, file name, and MIME type.
                 Returns error message string if creation fails.
        
        Raises:
            No exceptions are raised - all errors are caught and returned as error message strings.
        """
        if not PPTX_AVAILABLE:
            logger.error("python-pptx library not available")
            return "Error: python-pptx library is not installed. Please install it with: pip install python-pptx"
        
        try:
            logger.info("Starting PowerPoint presentation creation")
            logger.debug(f"Presentation title: {title}")
            logger.debug(f"Number of slides: {len(slides)}")
            logger.debug(f"Template: {template if template else 'default'}")
            logger.debug(f"Output format: {output_format}")
            
            # Create presentation
            try:
                if template:
                    logger.debug(f"Attempting to load template: {template}")
                    try:
                        prs = Presentation(template)
                        logger.info(f"Template loaded successfully: {template}")
                    except FileNotFoundError:
                        logger.warning(f"Template file not found: {template}, using default")
                        prs = Presentation()
                    except Exception as template_error:
                        logger.warning(f"Error loading template {template}: {str(template_error)}, using default")
                        prs = Presentation()
                else:
                    logger.debug("Creating presentation with default template")
                    prs = Presentation()
                    logger.debug("Default presentation created")
            except Exception as create_error:
                logger.error(f"Error creating presentation object: {str(create_error)}", exc_info=True)
                return f"Error: Failed to create presentation - {str(create_error)}"
            
            # Add title slide
            try:
                logger.debug("Adding title slide")
                title_slide_layout = prs.slide_layouts[0]  # Title slide layout
                slide = prs.slides.add_slide(title_slide_layout)
                title_shape = slide.shapes.title
                subtitle_shape = slide.placeholders[1]
                title_shape.text = title
                subtitle_shape.text = "Generated Presentation"
                logger.debug("Title slide added successfully")
            except Exception as title_slide_error:
                logger.error(f"Error adding title slide: {str(title_slide_error)}", exc_info=True)
                return f"Error: Failed to add title slide - {str(title_slide_error)}"
            
            # Add content slides
            slides_added = 0
            slides_failed = 0
            
            for slide_idx, slide_config in enumerate(slides, 1):
                try:
                    logger.debug(f"Processing slide {slide_idx}/{len(slides)}")
                    layout_type = slide_config.get("layout", "content")
                    logger.debug(f"Slide layout type: {layout_type}")
                    
                    # Select layout based on type
                    try:
                        if layout_type == "title":
                            slide_layout = prs.slide_layouts[0]  # Title slide
                        elif layout_type == "content":
                            slide_layout = prs.slide_layouts[1]  # Title and Content
                        elif layout_type == "two_content":
                            slide_layout = prs.slide_layouts[3]  # Two Content
                        elif layout_type == "blank":
                            slide_layout = prs.slide_layouts[6]  # Blank
                        else:
                            logger.warning(f"Unknown layout type '{layout_type}', using default content layout")
                            slide_layout = prs.slide_layouts[1]  # Default: Title and Content
                        logger.debug(f"Layout selected: {type(slide_layout).__name__}")
                    except Exception as layout_error:
                        logger.error(f"Error selecting layout for slide {slide_idx}: {str(layout_error)}", exc_info=True)
                        slides_failed += 1
                        continue
                    
                    # Create slide
                    try:
                        slide = prs.slides.add_slide(slide_layout)
                        logger.debug(f"Slide {slide_idx} created")
                    except Exception as slide_create_error:
                        logger.error(f"Error creating slide {slide_idx}: {str(slide_create_error)}", exc_info=True)
                        slides_failed += 1
                        continue
                    
                    # Set title if layout has title placeholder
                    try:
                        if slide.shapes.title:
                            slide_title = slide_config.get("title", "")
                            slide.shapes.title.text = slide_title
                            logger.debug(f"Slide {slide_idx} title set: {slide_title}")
                    except Exception as title_error:
                        logger.warning(f"Error setting title for slide {slide_idx}: {str(title_error)}")
                    
                    # Add content
                    try:
                        content = slide_config.get("content", [])
                        if content and len(slide.placeholders) > 1:
                            logger.debug(f"Adding {len(content)} content blocks to slide {slide_idx}")
                            content_placeholder = slide.placeholders[1]
                            tf = content_placeholder.text_frame
                            tf.text = ""
                            
                            for block_idx, block in enumerate(content, 1):
                                try:
                                    if isinstance(block, str):
                                        p = tf.add_paragraph()
                                        p.text = block
                                        p.level = 0
                                        logger.debug(f"Added text paragraph {block_idx} to slide {slide_idx}")
                                    elif isinstance(block, dict):
                                        text = block.get("text", "")
                                        level = block.get("level", 0)
                                        p = tf.add_paragraph()
                                        p.text = text
                                        p.level = level
                                        logger.debug(f"Added structured paragraph {block_idx} (level {level}) to slide {slide_idx}")
                                except Exception as block_error:
                                    logger.warning(f"Error adding content block {block_idx} to slide {slide_idx}: {str(block_error)}")
                                    continue
                            logger.debug(f"Content added to slide {slide_idx}")
                    except Exception as content_error:
                        logger.warning(f"Error adding content to slide {slide_idx}: {str(content_error)}")
                    
                    # Add images if provided
                    try:
                        images = slide_config.get("images", [])
                        if images:
                            logger.debug(f"Adding {len(images)} image(s) to slide {slide_idx}")
                            for img_idx, img_data in enumerate(images, 1):
                                try:
                                    logger.debug(f"Processing image {img_idx}/{len(images)} for slide {slide_idx}")
                                    img_bytes = base64.b64decode(img_data)
                                    img_io = io.BytesIO(img_bytes)
                                    # Add image to slide (simplified - can be enhanced)
                                    slide.shapes.add_picture(img_io, Inches(1), Inches(2), width=Inches(6))
                                    logger.debug(f"Image {img_idx} added to slide {slide_idx}")
                                except Exception as img_error:
                                    logger.warning(f"Error adding image {img_idx} to slide {slide_idx}: {str(img_error)}")
                                    continue
                    except Exception as images_error:
                        logger.warning(f"Error processing images for slide {slide_idx}: {str(images_error)}")
                    
                    slides_added += 1
                    logger.debug(f"Slide {slide_idx} completed successfully")
                    
                except Exception as slide_error:
                    logger.error(f"Unexpected error processing slide {slide_idx}: {str(slide_error)}", exc_info=True)
                    slides_failed += 1
                    continue
            
            logger.info(f"Slides processed: {slides_added} succeeded, {slides_failed} failed")
            
            # Save to bytes
            try:
                logger.debug("Saving presentation to bytes")
                output = io.BytesIO()
                prs.save(output)
                output.seek(0)
                file_bytes = output.read()
                logger.info(f"Presentation saved successfully (size: {len(file_bytes)} bytes)")
            except Exception as save_error:
                logger.error(f"Error saving presentation: {str(save_error)}", exc_info=True)
                return f"Error: Failed to save presentation - {str(save_error)}"
            
            # Encode or save to file
            try:
                if output_format == "base64":
                    logger.debug("Encoding presentation to base64")
                    encoded = _encode_file_output(file_bytes)
                    safe_title = title.replace(' ', '_').replace('/', '_').replace('\\', '_')
                    logger.info("Successfully created PowerPoint presentation")
                    return json.dumps({
                        "status": "success",
                        "file_data": encoded,
                        "file_name": f"{safe_title}.pptx",
                        "mime_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        "slides_added": slides_added,
                        "slides_failed": slides_failed
                    })
                else:
                    logger.debug(f"Saving presentation to path: {output_format}")
                    try:
                        with open(output_format, 'wb') as f:
                            f.write(file_bytes)
                        logger.info(f"Presentation saved to: {output_format}")
                        return json.dumps({
                            "status": "success",
                            "file_path": output_format,
                            "slides_added": slides_added,
                            "slides_failed": slides_failed
                        })
                    except PermissionError as perm_error:
                        logger.error(f"Permission denied writing to file: {output_format}")
                        return f"Error: Permission denied writing to file - {str(perm_error)}"
                    except Exception as file_write_error:
                        logger.error(f"Error writing file: {str(file_write_error)}", exc_info=True)
                        return f"Error: Failed to write file - {str(file_write_error)}"
            except Exception as output_error:
                logger.error(f"Error in output formatting: {str(output_error)}", exc_info=True)
                return f"Error: Failed to format output - {str(output_error)}"
            
        except ValueError as val_error:
            logger.error(f"Value error in PowerPoint creation: {str(val_error)}", exc_info=True)
            return f"Error: Invalid input - {str(val_error)}"
        except MemoryError as mem_error:
            logger.error(f"Memory error creating PowerPoint (file may be too large): {str(mem_error)}", exc_info=True)
            return f"Error: File too large to process - {str(mem_error)}"
        except Exception as e:
            logger.error(f"Unexpected error creating PowerPoint presentation: {str(e)}", exc_info=True)
            return f"Error creating PowerPoint presentation: {str(e)}"
    
    async def _arun(
        self,
        title: str,
        slides: List[Dict[str, Any]],
        template: Optional[str] = None,
        output_format: str = "base64"
    ) -> str:
        """
        Async version of _run for better performance in async contexts.
        
        This method provides an asynchronous interface for presentation creation,
        allowing it to be used efficiently in async agent execution flows. Currently
        delegates to the synchronous _run method, but structured for future async optimization.
        
        Args:
            title: Presentation title
            slides: List of slide configuration dictionaries
            template: Optional template file path
            output_format: Output format - "base64" or "file_path"
        
        Returns:
            str: JSON string with PowerPoint file, same as _run().
        
        Note:
            Currently implemented as a wrapper around _run() for compatibility.
            Future versions may implement true async file I/O for better performance.
        """
        return self._run(title, slides, template, output_format)


class AddSlideInput(BaseModel):
    """Input schema for AddSlideTool."""
    file_path: str = Field(..., description="Path to PowerPoint file or base64-encoded file data")
    slide_config: Dict[str, Any] = Field(..., description="Slide configuration dictionary")
    position: Optional[int] = Field(None, description="Position to insert slide (default: append at end)")
    output_format: str = Field("base64", description="Output format: 'base64' or 'file_path' (default: 'base64')")


class AddSlideTool(BaseTool):
    """
    Add a slide to an existing PowerPoint presentation.
    
    Slide configuration includes:
    - layout: Slide layout type
    - title: Slide title
    - content: List of content blocks
    - images: Optional list of base64-encoded images
    - charts: Optional list of chart configurations
    
    Use this tool when you need to:
    - Add slides to existing presentations
    - Update presentations with new content
    - Append analysis results to presentations
    """
    name: str = "add_slide_to_presentation"
    description: str = (
        "Add a slide to an existing PowerPoint presentation. "
        "Returns base64-encoded PowerPoint file with new slide added."
    )
    args_schema: Type[BaseModel] = AddSlideInput

    def _run(
        self,
        file_path: str,
        slide_config: Dict[str, Any],
        position: Optional[int] = None,
        output_format: str = "base64"
    ) -> str:
        """
        Add slide to PowerPoint presentation.
        
        This method performs the slide addition:
        1. Decodes file input (base64 or file path)
        2. Loads existing presentation using python-pptx
        3. Selects appropriate slide layout
        4. Creates new slide at specified position or appends
        5. Sets slide title and content
        6. Embeds images if provided
        7. Saves modified presentation and returns base64-encoded file
        
        Args:
            file_path: Path to PowerPoint file or base64-encoded file data.
                      Example: "presentation.pptx" or base64 string from file upload
            slide_config: Dictionary containing slide configuration:
                        - layout: Layout type (title, content, two_content, blank)
                        - title: Slide title
                        - content: List of content blocks (strings or dicts)
                        - images: Optional list of base64-encoded images
            position: Optional position to insert slide (0-based index). If None, appends at end.
                     Example: 0 (first slide), 3 (fourth slide), None (append)
            output_format: Output format - "base64" (default) or "file_path"
        
        Returns:
            str: JSON string with status, base64-encoded PowerPoint file with new slide, file name, and MIME type.
                 Returns error message string if addition fails.
        
        Raises:
            No exceptions are raised - all errors are caught and returned as error message strings.
        """
        if not PPTX_AVAILABLE:
            logger.error("python-pptx library not available")
            return "Error: python-pptx library is not installed. Please install it with: pip install python-pptx"
        
        try:
            logger.info("Starting slide addition to PowerPoint presentation")
            logger.debug(f"Position: {position if position is not None else 'append at end'}")
            logger.debug(f"Output format: {output_format}")
            logger.debug(f"Slide config layout: {slide_config.get('layout', 'content')}")
            
            # Decode file input
            try:
                logger.debug("Decoding file input")
                file_bytes, is_base64 = _decode_file_input(file_path)
                logger.debug(f"File decoded successfully (size: {len(file_bytes)} bytes, is_base64: {is_base64})")
                file_obj = io.BytesIO(file_bytes)
            except ValueError as decode_error:
                logger.error(f"Error decoding file input: {str(decode_error)}", exc_info=True)
                return f"Error: Failed to decode file input - {str(decode_error)}"
            except Exception as decode_error:
                logger.error(f"Unexpected error decoding file input: {str(decode_error)}", exc_info=True)
                return f"Error: Failed to decode file input - {str(decode_error)}"
            
            # Load presentation
            try:
                logger.debug("Loading PowerPoint presentation")
                prs = Presentation(file_obj)
                initial_slide_count = len(prs.slides)
                logger.info(f"Presentation loaded successfully ({initial_slide_count} slides)")
            except Exception as load_error:
                logger.error(f"Error loading presentation: {str(load_error)}", exc_info=True)
                return f"Error: Failed to load presentation - {str(load_error)}"
            
            # Select layout
            try:
                layout_type = slide_config.get("layout", "content")
                logger.debug(f"Selecting layout type: {layout_type}")
                
                if layout_type == "title":
                    slide_layout = prs.slide_layouts[0]
                elif layout_type == "content":
                    slide_layout = prs.slide_layouts[1]
                elif layout_type == "two_content":
                    slide_layout = prs.slide_layouts[3]
                elif layout_type == "blank":
                    slide_layout = prs.slide_layouts[6]
                else:
                    logger.warning(f"Unknown layout type '{layout_type}', using default content layout")
                    slide_layout = prs.slide_layouts[1]
                logger.debug(f"Layout selected: {type(slide_layout).__name__}")
            except Exception as layout_error:
                logger.error(f"Error selecting layout: {str(layout_error)}", exc_info=True)
                return f"Error: Failed to select layout - {str(layout_error)}"
            
            # Add slide at position or append
            try:
                if position is not None:
                    if 0 <= position < len(prs.slides):
                        logger.debug(f"Inserting slide at position {position}")
                        slide = prs.slides.add_slide(slide_layout, prs.slides._sldIdLst[position])
                        logger.debug(f"Slide inserted at position {position}")
                    else:
                        logger.warning(f"Position {position} out of range (0-{len(prs.slides)-1}), appending at end")
                        slide = prs.slides.add_slide(slide_layout)
                        logger.debug("Slide appended at end (position was out of range)")
                else:
                    logger.debug("Appending slide at end")
                    slide = prs.slides.add_slide(slide_layout)
                    logger.debug("Slide appended successfully")
            except Exception as slide_add_error:
                logger.error(f"Error adding slide: {str(slide_add_error)}", exc_info=True)
                return f"Error: Failed to add slide - {str(slide_add_error)}"
            
            # Set title
            try:
                if slide.shapes.title:
                    slide_title = slide_config.get("title", "")
                    slide.shapes.title.text = slide_title
                    logger.debug(f"Slide title set: {slide_title}")
                else:
                    logger.debug("Slide layout does not have title placeholder")
            except Exception as title_error:
                logger.warning(f"Error setting slide title: {str(title_error)}")
            
            # Add content
            try:
                content = slide_config.get("content", [])
                if content and len(slide.placeholders) > 1:
                    logger.debug(f"Adding {len(content)} content blocks to slide")
                    content_placeholder = slide.placeholders[1]
                    tf = content_placeholder.text_frame
                    tf.text = ""
                    
                    for block_idx, block in enumerate(content, 1):
                        try:
                            if isinstance(block, str):
                                p = tf.add_paragraph()
                                p.text = block
                                p.level = 0
                                logger.debug(f"Added text paragraph {block_idx}")
                            elif isinstance(block, dict):
                                text = block.get("text", "")
                                level = block.get("level", 0)
                                p = tf.add_paragraph()
                                p.text = text
                                p.level = level
                                logger.debug(f"Added structured paragraph {block_idx} (level {level})")
                        except Exception as block_error:
                            logger.warning(f"Error adding content block {block_idx}: {str(block_error)}")
                            continue
                    logger.debug("Content added to slide")
                else:
                    logger.debug("No content to add or slide has insufficient placeholders")
            except Exception as content_error:
                logger.warning(f"Error adding content to slide: {str(content_error)}")
            
            # Add images if provided
            try:
                images = slide_config.get("images", [])
                if images:
                    logger.debug(f"Adding {len(images)} image(s) to slide")
                    for img_idx, img_data in enumerate(images, 1):
                        try:
                            logger.debug(f"Processing image {img_idx}/{len(images)}")
                            img_bytes = base64.b64decode(img_data)
                            img_io = io.BytesIO(img_bytes)
                            slide.shapes.add_picture(img_io, Inches(1), Inches(2), width=Inches(6))
                            logger.debug(f"Image {img_idx} added to slide")
                        except Exception as img_error:
                            logger.warning(f"Error adding image {img_idx}: {str(img_error)}")
                            continue
            except Exception as images_error:
                logger.warning(f"Error processing images: {str(images_error)}")
            
            # Save to bytes
            try:
                logger.debug("Saving presentation to bytes")
                output = io.BytesIO()
                prs.save(output)
                output.seek(0)
                file_bytes = output.read()
                final_slide_count = len(prs.slides)
                logger.info(f"Presentation saved successfully (size: {len(file_bytes)} bytes, slides: {final_slide_count})")
            except Exception as save_error:
                logger.error(f"Error saving presentation: {str(save_error)}", exc_info=True)
                return f"Error: Failed to save presentation - {str(save_error)}"
            
            # Encode or save to file
            try:
                if output_format == "base64":
                    logger.debug("Encoding presentation to base64")
                    encoded = _encode_file_output(file_bytes)
                    logger.info("Successfully added slide to presentation")
                    return json.dumps({
                        "status": "success",
                        "file_data": encoded,
                        "file_name": "presentation_with_slide.pptx",
                        "mime_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    })
                else:
                    logger.debug(f"Saving presentation to path: {output_format}")
                    try:
                        with open(output_format, 'wb') as f:
                            f.write(file_bytes)
                        logger.info(f"Presentation saved to: {output_format}")
                        return json.dumps({"status": "success", "file_path": output_format})
                    except PermissionError as perm_error:
                        logger.error(f"Permission denied writing to file: {output_format}")
                        return f"Error: Permission denied writing to file - {str(perm_error)}"
                    except Exception as file_write_error:
                        logger.error(f"Error writing file: {str(file_write_error)}", exc_info=True)
                        return f"Error: Failed to write file - {str(file_write_error)}"
            except Exception as output_error:
                logger.error(f"Error in output formatting: {str(output_error)}", exc_info=True)
                return f"Error: Failed to format output - {str(output_error)}"
            
        except ValueError as val_error:
            logger.error(f"Value error adding slide: {str(val_error)}", exc_info=True)
            return f"Error: Invalid input - {str(val_error)}"
        except MemoryError as mem_error:
            logger.error(f"Memory error adding slide (file may be too large): {str(mem_error)}", exc_info=True)
            return f"Error: File too large to process - {str(mem_error)}"
        except Exception as e:
            logger.error(f"Unexpected error adding slide to presentation: {str(e)}", exc_info=True)
            return f"Error adding slide to presentation: {str(e)}"
    
    async def _arun(
        self,
        file_path: str,
        slide_config: Dict[str, Any],
        position: Optional[int] = None,
        output_format: str = "base64"
    ) -> str:
        """
        Async version of _run for better performance in async contexts.
        
        This method provides an asynchronous interface for slide addition,
        allowing it to be used efficiently in async agent execution flows. Currently
        delegates to the synchronous _run method, but structured for future async optimization.
        
        Args:
            file_path: Path to PowerPoint file or base64-encoded file data
            slide_config: Dictionary containing slide configuration
            position: Optional position to insert slide
            output_format: Output format - "base64" or "file_path"
        
        Returns:
            str: JSON string with PowerPoint file containing new slide, same as _run().
        
        Note:
            Currently implemented as a wrapper around _run() for compatibility.
            Future versions may implement true async file I/O for better performance.
        """
        return self._run(file_path, slide_config, position, output_format)


class AddChartToSlideInput(BaseModel):
    """Input schema for AddChartToSlideTool."""
    file_path: str = Field(..., description="Path to PowerPoint file or base64-encoded file data")
    slide_index: int = Field(..., description="Slide number (0-based index)")
    chart_config: Dict[str, Any] = Field(..., description="Chart configuration dictionary")
    output_format: str = Field("base64", description="Output format: 'base64' or 'file_path' (default: 'base64')")


class AddChartToSlideTool(BaseTool):
    """
    Add a chart to a PowerPoint slide.
    
    This tool adds charts to PowerPoint slides by generating chart images and embedding them.
    It supports various chart types (bar, line, pie, scatter) and allows positioning and sizing
    control. Charts are generated using matplotlib or plotly and embedded as images in the slide.
    
    **When to use this tool:**
    - Adding visualizations to PowerPoint slides
    - Embedding charts in presentations
    - Creating data-driven presentations
    - Adding analysis charts to report slides
    - Visualizing data in presentation format
    
    **When NOT to use this tool:**
    - For creating standalone charts (use visualization tools)
    - When you need Excel-style embedded charts (not supported)
    - For interactive charts (charts are static images)
    - When chart data is not available
    
    **Input requirements:**
    - Must provide a valid PowerPoint file (.pptx format) as base64 or file path
    - Slide index must be valid (0-based, within presentation bounds)
    - Chart configuration must include chart type and data
    - Position and size are optional (defaults provided)
    
    **Output:**
    - Returns JSON with status, base64-encoded PowerPoint file with chart, file name, and MIME type
    - Chart is embedded as image in the specified slide
    - Returns error message if chart addition fails
    
    **Supported Chart Types:**
    - bar: Bar chart
    - line: Line chart
    - pie: Pie chart
    - scatter: Scatter plot
    
    **Limitations:**
    - Charts are embedded as images (not native PowerPoint charts)
    - Requires matplotlib/plotly integration (currently placeholder)
    - Chart customization options are limited
    - Cannot edit charts in PowerPoint (they're images)
    
    **Example use cases:**
    1. "Add a bar chart to slide 2 showing the analysis results"
    2. "Embed a pie chart in the summary slide"
    3. "Add a line chart to slide 5 visualizing the trends"
    4. "Create a scatter plot on slide 3 with the correlation data"
    
    **Configuration:**
    Requires python-pptx and matplotlib/plotly libraries. The tool generates chart images
    and embeds them in slides. Full implementation requires matplotlib/plotly integration.
    Currently returns a placeholder message indicating the requirement.
    """
    name: str = "add_chart_to_slide"
    description: str = (
        "Add a chart to a PowerPoint slide. "
        "Supports various chart types (bar, line, pie, scatter). "
        "Returns base64-encoded PowerPoint file with chart added. "
        "IMPORTANT: File can be base64-encoded string or file path. Slide index is 0-based. "
        "NOTE: Full chart embedding requires matplotlib/plotly integration (currently placeholder). "
        "Best for: adding visualizations to slides, embedding charts in presentations. "
        "NOT suitable for: standalone charts, interactive visualizations, or Excel-style embedded charts."
    )
    args_schema: Type[BaseModel] = AddChartToSlideInput

    def _run(
        self,
        file_path: str,
        slide_index: int,
        chart_config: Dict[str, Any],
        output_format: str = "base64"
    ) -> str:
        """
        Add chart to PowerPoint slide.
        
        This method performs the chart addition (placeholder implementation):
        1. Decodes file input (base64 or file path)
        2. Loads existing presentation
        3. Validates slide index
        4. Generates chart image using matplotlib/plotly (requires integration)
        5. Embeds chart image in slide at specified position
        6. Saves modified presentation and returns base64-encoded file
        
        Args:
            file_path: Path to PowerPoint file or base64-encoded file data.
                      Example: "presentation.pptx" or base64 string from file upload
            slide_index: Slide number (0-based index) where chart should be added.
                        Example: 0 (first slide), 2 (third slide)
            chart_config: Dictionary containing chart configuration:
                        - chart_type: Type of chart (bar, line, pie, scatter)
                        - data: Chart data (list of dicts or 2D array)
                        - position: Chart position {"left": 1, "top": 2} in inches
                        - size: Chart size {"width": 6, "height": 4} in inches
                        - title: Optional chart title
            output_format: Output format - "base64" (default) or "file_path"
        
        Returns:
            str: JSON string with status, base64-encoded PowerPoint file with chart, file name, and MIME type.
                 Returns placeholder message indicating matplotlib/plotly integration requirement.
                 Returns error message string if addition fails.
        
        Raises:
            No exceptions are raised - all errors are caught and returned as error message strings.
        
        Note:
            Full implementation requires matplotlib/plotly integration to generate chart images.
            Currently returns a placeholder message. To complete: generate chart image, embed in slide.
        """
        if not PPTX_AVAILABLE:
            logger.error("python-pptx library not available")
            return "Error: python-pptx library is not installed. Please install it with: pip install python-pptx"
        
        try:
            logger.info(f"Starting chart addition to slide {slide_index}")
            logger.debug(f"Output format: {output_format}")
            logger.debug(f"Chart config: type={chart_config.get('chart_type', 'bar')}, has_data={bool(chart_config.get('data'))}")
            
            # Decode file input
            try:
                logger.debug("Decoding file input")
                file_bytes, is_base64 = _decode_file_input(file_path)
                logger.debug(f"File decoded successfully (size: {len(file_bytes)} bytes, is_base64: {is_base64})")
                file_obj = io.BytesIO(file_bytes)
            except ValueError as decode_error:
                logger.error(f"Error decoding file input: {str(decode_error)}", exc_info=True)
                return f"Error: Failed to decode file input - {str(decode_error)}"
            except Exception as decode_error:
                logger.error(f"Unexpected error decoding file input: {str(decode_error)}", exc_info=True)
                return f"Error: Failed to decode file input - {str(decode_error)}"
            
            # Load presentation
            try:
                logger.debug("Loading PowerPoint presentation")
                prs = Presentation(file_obj)
                slide_count = len(prs.slides)
                logger.info(f"Presentation loaded successfully ({slide_count} slides)")
            except Exception as load_error:
                logger.error(f"Error loading presentation: {str(load_error)}", exc_info=True)
                return f"Error: Failed to load presentation - {str(load_error)}"
            
            # Validate slide index
            try:
                slide_count = len(prs.slides)
                if slide_index < 0:
                    logger.error(f"Invalid slide index: {slide_index} (negative)")
                    return f"Error: Invalid slide index {slide_index}. Slide index must be 0 or greater."
                if slide_index >= slide_count:
                    logger.error(f"Invalid slide index: {slide_index} (presentation has {slide_count} slides)")
                    return f"Error: Invalid slide index {slide_index}. Presentation has {slide_count} slides (0-{slide_count-1})."
                logger.debug(f"Slide index {slide_index} validated (presentation has {slide_count} slides)")
            except Exception as validation_error:
                logger.error(f"Error validating slide index: {str(validation_error)}", exc_info=True)
                return f"Error: Failed to validate slide index - {str(validation_error)}"
            
            # Get slide
            try:
                slide = prs.slides[slide_index]
                logger.debug(f"Slide {slide_index} retrieved successfully")
            except IndexError as idx_error:
                logger.error(f"Slide index {slide_index} out of bounds", exc_info=True)
                return f"Error: Slide index {slide_index} is out of bounds."
            except Exception as slide_error:
                logger.error(f"Error retrieving slide {slide_index}: {str(slide_error)}", exc_info=True)
                return f"Error: Failed to retrieve slide - {str(slide_error)}"
            
            # Extract chart configuration
            try:
                chart_type = chart_config.get("chart_type", "bar")
                data = chart_config.get("data", [])
                position = chart_config.get("position", {"left": 1, "top": 2})
                size = chart_config.get("size", {"width": 6, "height": 4})
                logger.debug(f"Chart configuration extracted: type={chart_type}, data_points={len(data) if isinstance(data, list) else 'N/A'}, position={position}, size={size}")
            except Exception as config_error:
                logger.warning(f"Error extracting chart configuration: {str(config_error)}, using defaults")
                chart_type = "bar"
                data = []
                position = {"left": 1, "top": 2}
                size = {"width": 6, "height": 4}
            
            # Note: python-pptx doesn't directly support charts
            # We need to generate chart as image and embed it
            # This is a simplified version - full implementation would:
            # 1. Generate chart using matplotlib/plotly
            # 2. Save as image
            # 3. Embed image in slide
            
            logger.warning("Chart embedding requires matplotlib/plotly integration - returning placeholder")
            logger.debug("Skipping chart generation (integration pending)")
            
            # Save to bytes
            try:
                logger.debug("Saving presentation to bytes")
                output = io.BytesIO()
                prs.save(output)
                output.seek(0)
                file_bytes = output.read()
                logger.info(f"Presentation saved successfully (size: {len(file_bytes)} bytes)")
            except Exception as save_error:
                logger.error(f"Error saving presentation: {str(save_error)}", exc_info=True)
                return f"Error: Failed to save presentation - {str(save_error)}"
            
            # Encode or save to file
            try:
                if output_format == "base64":
                    logger.debug("Encoding presentation to base64")
                    encoded = _encode_file_output(file_bytes)
                    logger.info("Successfully processed presentation (chart embedding pending)")
                    return json.dumps({
                        "status": "success",
                        "file_data": encoded,
                        "file_name": "presentation_with_chart.pptx",
                        "mime_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        "note": "Chart embedding requires matplotlib/plotly - implementation pending"
                    })
                else:
                    logger.debug(f"Saving presentation to path: {output_format}")
                    try:
                        with open(output_format, 'wb') as f:
                            f.write(file_bytes)
                        logger.info(f"Presentation saved to: {output_format}")
                        return json.dumps({
                            "status": "success",
                            "file_path": output_format,
                            "note": "Chart embedding requires matplotlib/plotly - implementation pending"
                        })
                    except PermissionError as perm_error:
                        logger.error(f"Permission denied writing to file: {output_format}")
                        return f"Error: Permission denied writing to file - {str(perm_error)}"
                    except Exception as file_write_error:
                        logger.error(f"Error writing file: {str(file_write_error)}", exc_info=True)
                        return f"Error: Failed to write file - {str(file_write_error)}"
            except Exception as output_error:
                logger.error(f"Error in output formatting: {str(output_error)}", exc_info=True)
                return f"Error: Failed to format output - {str(output_error)}"
            
        except ValueError as val_error:
            logger.error(f"Value error adding chart: {str(val_error)}", exc_info=True)
            return f"Error: Invalid input - {str(val_error)}"
        except MemoryError as mem_error:
            logger.error(f"Memory error adding chart (file may be too large): {str(mem_error)}", exc_info=True)
            return f"Error: File too large to process - {str(mem_error)}"
        except Exception as e:
            logger.error(f"Unexpected error adding chart to slide: {str(e)}", exc_info=True)
            return f"Error adding chart to slide: {str(e)}"
    
    async def _arun(
        self,
        file_path: str,
        slide_index: int,
        chart_config: Dict[str, Any],
        output_format: str = "base64"
    ) -> str:
        """
        Async version of _run for better performance in async contexts.
        
        This method provides an asynchronous interface for chart addition,
        allowing it to be used efficiently in async agent execution flows. Currently
        delegates to the synchronous _run method, but structured for future async optimization.
        
        Args:
            file_path: Path to PowerPoint file or base64-encoded file data
            slide_index: Slide number (0-based index)
            chart_config: Dictionary containing chart configuration
            output_format: Output format - "base64" or "file_path"
        
        Returns:
            str: JSON string with PowerPoint file containing chart, same as _run().
        
        Note:
            Currently implemented as a wrapper around _run() for compatibility.
            Future versions may implement true async file I/O for better performance.
        """
        return self._run(file_path, slide_index, chart_config, output_format)

